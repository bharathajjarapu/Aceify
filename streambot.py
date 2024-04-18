import os
import io
import pip
import PIL
import docx
import pptx
import logging
import subprocess
import pandas as pd
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import google.generativeai as gen_ai
from langchain_community.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from telegram import Update
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes
import asyncio

load_dotenv()
os.getenv("GOOGLE_API_KEY")
gen_ai.configure(api_key=os.environ['GOOGLE_API_KEY'])
st.set_page_config(page_title="Aceify", page_icon=":book:")

ADMIN = 5882873021
USER_IDS = 1111111111, 0000000000

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "messages" not in st.session_state:
    st.session_state.messages = []

logging.basicConfig(format='%(asctime)s - %(name)s - %(levelname)s - %(message)s', level=logging.INFO)
logger = logging.getLogger(__name__)

user_vector_stores = {}

def extract_text_from_file(file_bytes, file_extension, ocr=False):
    text = ""

    if file_extension == ".pdf":
        pdf_io = io.BytesIO(file_bytes)
        pdf_reader = PdfReader(pdf_io)
        for page in pdf_reader.pages:
            text += page.extract_text()

    elif file_extension in [".png", ".jpg", ".jpeg", ".webp"]:
        image_file = io.BytesIO(file_bytes)
        image = PIL.Image.open(image_file)
        vision_model = gen_ai.GenerativeModel('gemini-pro-vision')
        response = vision_model.generate_content(["Explain the picture and Extract Text from picture?",image])
        text = response.text

    elif file_extension == ".docx":
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
    
    elif file_extension == ".pptx":
        presentation = pptx.Presentation(io.BytesIO(file_bytes))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file_extension == ".csv":
        data = pd.read_csv(io.BytesIO(file_bytes))
        text = data.to_string()

    return text

def split_text_into_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    text_chunks = text_splitter.split_text(text)
    return text_chunks

def create_vector_store(user_id, chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    vector_store.save_local(f"faiss_index_{user_id}")
    user_vector_stores[user_id] = vector_store

def create_vectors(chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    st.session_state.vector_store = vector_store

def setup_conversational_chain():
    prompt_template = """
    You are an AI RAG ChatBot, you help in extracting insights and information from the given context for the question, you must follow the given rules
    Rule 1: Answer the question in a short, precise, detailed, friendly and engaging way, drawing from the provided context if possible. 
    Rule 2: If the question is not directly related to the context, provide a thoughtful and relevant response based on your general knowledge.
    Rule 3: Make sure the response is in common English and appears human-like, 
    Rule 4: Avoid any patterns or language in response that could be identified as AI-generated by AI text detectors.
    Rule 5: Always use LaTeX when formulaes are involved

    Context:{context}

    Question:{question}

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", client=gen_ai, temperature=0.7)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(llm=model, chain_type="stuff", prompt=prompt)
    return chain

def get_user_response(user_id, user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.load_local(f"faiss_index_{user_id}", embeddings, allow_dangerous_deserialization=True)
    relevant_docs = vector_store.similarity_search(user_question)

    conversational_chain = setup_conversational_chain()
    response = conversational_chain({"input_documents": relevant_docs, "question": user_question}, return_only_outputs=True)

    return response

def get_response(user_question):
    relevant_docs = st.session_state.vector_store.similarity_search(user_question)
    conversational_chain = setup_conversational_chain()
    response = conversational_chain({"input_documents": relevant_docs, "question": user_question}, return_only_outputs=True)
    return response

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id
    if int(user_id) == ADMIN:
        await update.message.reply_markdown_v2(
            f"Hi {update.effective_user.mention_markdown_v2()}!\nWelcome Back Sir"
        )
    elif user_id in USER_IDS:
        await update.message.reply_markdown_v2(
            f"Hi {update.effective_user.mention_markdown_v2()}!\nHi This is Aceify."
        )
    else:
        await update.message.reply_text("Sorry, you don't have access to this bot.")

async def clear(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    import os
    import shutil

    directory = os.path.abspath(os.path.dirname(__file__))
    for filename in os.listdir(directory):
        if filename.startswith("faiss_index_"):
            file_path = os.path.join(directory, filename)
            os.remove(file_path)
    user_vector_stores.clear()
    await update.message.reply_text("All FAISS indexes have been deleted.")

async def process_file(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id

    if update.message.document:
        file = await update.message.document.get_file()
    elif update.message.photo:
        file = await update.message.photo[-1].get_file()
    else:
        await update.message.reply_text("Please upload a document or photo.")
        return

    file_bytes = await file.download_as_bytearray()
    file_extension = os.path.splitext(file.file_path)[1].lower()

    await update.message.reply_text(f"Processing your file...")

    file_text = extract_text_from_file(file_bytes, file_extension, ocr=True)
    text_chunks = split_text_into_chunks(file_text)
    create_vector_store(user_id, text_chunks)

    await update.message.reply_text(f"File processed successfully.")

async def askque(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id
    user_question = update.message.text

    if user_id not in user_vector_stores:
        await update.message.reply_text("Please upload your Files ?!")
        return

    response = get_user_response(user_id, user_question)
    full_response = ''.join(response['output_text'])
    await update.message.reply_text(full_response)

def main():
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)

    application = ApplicationBuilder().token(os.environ['TELEGRAM_BOT_TOKEN']).build()
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("clear", clear))
    application.add_handler(MessageHandler(filters.Document.ALL, process_file))
    application.add_handler(MessageHandler(filters.PHOTO, process_file))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, askque))

    st.title("Aceify")
    st.write("Welcome to Aceify! Upload your syllabus PDFs and ask questions.")

    uploaded_file = st.file_uploader("Upload your Documents (PDF, Word, PowerPoint, CSV, or spreadsheet)", type=["pdf", "docx", "pptx", "csv", "xls", "xlsx", "png", "jpg", "jpeg"])
    ocr = st.checkbox("Enable OCR for PDF files")

    if st.button("Clear Conversation"):
        st.session_state.messages.clear()

    if uploaded_file is not None:
        file_bytes = uploaded_file.read()
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()

        if file_extension in [".png", ".jpg", ".jpeg"]:
            with st.expander("Display Image"):
                st.image(uploaded_file)

        st.write(f"Processing your {uploaded_file.name} file...")
        file_text = extract_text_from_file(file_bytes, file_extension, ocr)
        text_chunks = split_text_into_chunks(file_text)
        create_vectors(text_chunks)
        st.write(f"{uploaded_file.name} file processed successfully!")

    if st.session_state.vector_store is not None:
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])

        if prompt := st.chat_input():
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = get_response(prompt)
                    full_response = ''.join(response['output_text'])
                    st.write(full_response)
                    message = {"role": "assistant", "content": full_response}
                    st.session_state.messages.append(message)

    async def run_telegram_bot():
        await application.run_polling()

    loop.create_task(run_telegram_bot())
    loop.run_forever()

if __name__ == '__main__':
    main()
