import os
import streamlit as st
import io
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain_community.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
import google.generativeai as gen_ai
import docx
import pptx
import pandas as pd

load_dotenv()
os.getenv("GOOGLE_API_KEY")
gen_ai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
st.set_page_config(page_title="Aceify", page_icon=":book:")

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "messages" not in st.session_state:
    st.session_state.messages = []

gen_ai.configure(api_key=os.environ["GOOGLE_API_KEY"])

def extract_text_from_file(file_bytes, file_extension):
    text = ""
    if file_extension == ".pdf":
        pdf_io = io.BytesIO(file_bytes)
        pdf_reader = PdfReader(pdf_io)
        for page in pdf_reader.pages:
            text += page.extract_text()
    elif file_extension == ".docx":
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_extension == ".pptx":
        presentation = pptx.Presentation(io.BytesIO(file_bytes))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_extension == ".csv":
        data = pd.read_csv(io.BytesIO(file_bytes))
        text = data.to_string()
    return text

def split_text_into_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    text_chunks = text_splitter.split_text(text)
    return text_chunks

def create_vector_store(chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    st.session_state.vector_store = vector_store

def setup_conversational_chain():
    prompt_template = """
    Answer the question in a short, precise, detailed, friendly and engaging way, drawing from the provided context if possible. If the question is not directly related to the context, provide a thoughtful and relevant response based on your general knowledge.

    Context:{context}

    Question:{question}

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", client=gen_ai, temperature=0.7)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(llm=model, chain_type="stuff", prompt=prompt)
    return chain

def get_response(user_question):
    relevant_docs = st.session_state.vector_store.similarity_search(user_question)
    conversational_chain = setup_conversational_chain()
    response = conversational_chain({"input_documents": relevant_docs, "question": user_question}, return_only_outputs=True)
    return response

def main():
    st.title("Aceify Bot")
    st.write("Welcome to Aceify! I'm your friendly files assistant. Upload your study materials, and I'll process them. Then, you can ask me questions about the content.")

    uploaded_file = st.file_uploader("Upload your study materials (PDF, Word, PowerPoint, CSV, or spreadsheet)", type=["pdf", "docx", "pptx", "csv", "xls", "xlsx"])

    if uploaded_file is not None:
        file_bytes = uploaded_file.read()
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        st.write(f"Processing your {file_extension} file...")

        file_text = extract_text_from_file(file_bytes, file_extension)
        text_chunks = split_text_into_chunks(file_text)
        create_vector_store(text_chunks)

        st.write(f"{file_extension} file processed successfully!")

    if st.session_state.vector_store is not None:
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])

        if prompt := st.chat_input():
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = get_response(prompt)
                    full_response = ''.join(response['output_text'])
                    st.write(full_response)
                    message = {"role": "assistant", "content": full_response}
                    st.session_state.messages.append(message)

if __name__ == "__main__":
    main()