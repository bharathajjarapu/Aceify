import os
import io
import pip
import PIL
import docx
import pptx
import pandas as pd
import logging
from dotenv import load_dotenv
import google.generativeai as gen_ai
from PyPDF2 import PdfReader
from langchain_community.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from telegram import Update
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes

load_dotenv()
os.getenv("GOOGLE_API_KEY")
gen_ai.configure(api_key=os.environ['GOOGLE_API_KEY'])

logging.basicConfig(format='%(asctime)s - %(name)s - %(levelname)s - %(message)s', level=logging.INFO)
logger = logging.getLogger(__name__)

user_vector_stores = {}

def extract_text_from_file(file_bytes, file_extension, ocr=False):
    text = ""

    if file_extension == ".pdf":
        pdf_io = io.BytesIO(file_bytes)
        pdf_reader = PdfReader(pdf_io)
        if ocr:
            for page in pdf_reader.pages:
                images = page.images
                for image in images:
                    image_bytes = image.data
                    image_file = io.BytesIO(image_bytes)
                    image = PIL.Image.open(image_file)
                    text += pytesseract.image_to_string(image) + "\n"
        else:
            for page in pdf_reader.pages:
                text += page.extract_text()

    elif file_extension in [".png", ".jpg", ".jpeg", ".webp"]:
        image_file = io.BytesIO(file_bytes)
        image = PIL.Image.open(image_file)
        vision_model = gen_ai.GenerativeModel('gemini-pro-vision')
        response = vision_model.generate_content(["Explain the picture and Extract Text from picture?", image])
        text = response.text

    elif file_extension == ".docx":
        doc = docx.Document(io.BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"

    elif file_extension == ".pptx":
        presentation = pptx.Presentation(io.BytesIO(file_bytes))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file_extension == ".csv":
        data = pd.read_csv(io.BytesIO(file_bytes))
        text = data.to_string()

    return text

def split_text_into_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    text_chunks = text_splitter.split_text(text)
    return text_chunks

def create_vector_store(user_id, chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(chunks, embedding=embeddings)
    vector_store.save_local(f"faiss_index_{user_id}")
    user_vector_stores[user_id] = vector_store

def setup_conversational_chain():
    prompt_template = """
    You are an AI RAG ChatBot, you help in extracting insights and information from the given context for the question, you must follow the given rules
    Rule 1: Answer the question in a short, precise, detailed, friendly and engaging way, drawing from the provided context if possible. 
    Rule 2: If the question is not directly related to the context, provide a thoughtful and relevant response based on your general knowledge.
    Rule 3: Make sure the response is in common English and appears human-like, 
    Rule 4: Avoid any patterns or language in response that could be identified as AI-generated by AI text detectors.
    Rule 5: Always use LaTeX when formulaes are involved

    Context:{context}

    Question:{question}

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", client=gen_ai, temperature=0.7)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(llm=model, chain_type="stuff", prompt=prompt)
    return chain

def get_user_response(user_id, user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.load_local(f"faiss_index_{user_id}", embeddings, allow_dangerous_deserialization=True)
    relevant_docs = vector_store.similarity_search(user_question)

    conversational_chain = setup_conversational_chain()
    response = conversational_chain({"input_documents": relevant_docs, "question": user_question}, return_only_outputs=True)

    return response

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user = update.effective_user
    await update.message.reply_markdown_v2(
        f"Hi {user.mention_markdown_v2()}\!",
        f"Welcome to Aceify! I'm your friendly study assistant. Send me your syllabus PDFs, and I'll process them. Then, you can ask me questions about the content or anything else you'd like to discuss. Let's make learning fun!",
    )

async def clear(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    import os
    import shutil

    directory = os.path.abspath(os.path.dirname(__file__))
    for filename in os.listdir(directory):
        if filename.startswith("faiss_index_"):
            file_path = os.path.join(directory, filename)
            os.remove(file_path)
    user_vector_stores.clear()
    await update.message.reply_text("All FAISS indexes have been deleted.")

async def process_file(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id
    file = await update.message.document.get_file()
    file_bytes = await file.download_as_bytearray()
    file_extension = os.path.splitext(file.file_path)[1].lower()

    await update.message.reply_text(f"Processing your file...")

    file_text = extract_text_from_file(file_bytes, file_extension, ocr=True)
    text_chunks = split_text_into_chunks(file_text)
    create_vector_store(user_id, text_chunks)

    await update.message.reply_text(f"File processed successfully .")

async def askque(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id
    user_question = update.message.text

    if user_id not in user_vector_stores:
        await update.message.reply_text("Please upload your syllabus PDFs first ?!")
        return

    response = get_user_response(user_id, user_question)
    full_response = ''.join(response['output_text'])
    await update.message.reply_text(full_response)

def main() -> None:
    application = ApplicationBuilder().token(os.environ['TELEGRAM_BOT_TOKEN']).build()
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("clear", clear))
    application.add_handler(MessageHandler(filters.Document.ALL, process_file))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, askque))
    application.run_polling()
    
if __name__ == '__main__':
    main()